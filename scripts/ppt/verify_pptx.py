from pptx import Presentation

prs = Presentation(r'E:\我的桌面\入党积极分子思想汇报.pptx')
print(f"Total slides: {len(prs.slides)}\n")
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    first_texts = texts[:2] if texts else ['(empty)']
    print(f"Slide {i}: {' | '.join(first_texts)}")
